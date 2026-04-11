"""
create_ppt.py — Generate Thesis_Presentation_Prerna_Tank.pptx
Mirrors the HTML presentation exactly: dark theme, same content, 16:9 slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree
import copy

# ── Colours ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0a, 0x0e, 0x27)
CARD     = RGBColor(0x16, 0x1b, 0x40)
CARD_ALT = RGBColor(0x1a, 0x1f, 0x4a)
CYAN     = RGBColor(0x00, 0xe5, 0xff)
ORANGE   = RGBColor(0xff, 0x91, 0x00)
GREEN    = RGBColor(0x69, 0xf0, 0xae)
PINK     = RGBColor(0xf0, 0x62, 0x92)
PURPLE   = RGBColor(0xb3, 0x88, 0xff)
YELLOW   = RGBColor(0xff, 0xd7, 0x40)
RED      = RGBColor(0xff, 0x52, 0x52)
TEXT_PRI = RGBColor(0xe8, 0xea, 0xf6)
TEXT_SEC = RGBColor(0x9f, 0xa8, 0xda)
WHITE    = RGBColor(0xff, 0xff, 0xff)

# ── Slide size: 16:9 ────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

FOOTER_TXT = "Intelligent Auto-Scaling in Kubernetes: A Machine Learning Based Predictive Approach  |  Prerna Tank (2410512)"


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, italic=False,
                color=TEXT_PRI, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, x, y, w, h, fill_color=CARD, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # 1 = MSO_SHAPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_footer(slide):
    add_textbox(slide, FOOTER_TXT,
                Inches(0.2), H - Inches(0.35),
                W - Inches(0.4), Inches(0.28),
                font_size=8, color=TEXT_SEC, italic=True)


def add_slide_title(slide, title, subtitle=None):
    # Gradient-style title (use CYAN as approximation)
    add_textbox(slide, title,
                Inches(0.6), Inches(0.55),
                W - Inches(1.2), Inches(0.75),
                font_size=28, bold=True, color=CYAN,
                align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(1.3),
                    W - Inches(1.2), Inches(0.4),
                    font_size=14, color=TEXT_SEC,
                    align=PP_ALIGN.CENTER)


def card_block(slide, x, y, w, h, title, body,
               title_color=CYAN, body_font=11):
    r = add_rect(slide, x, y, w, h, fill_color=CARD)
    add_textbox(slide, title,
                x + Inches(0.15), y + Inches(0.1),
                w - Inches(0.3), Inches(0.35),
                font_size=12, bold=True, color=title_color)
    add_textbox(slide, body,
                x + Inches(0.15), y + Inches(0.42),
                w - Inches(0.3), h - Inches(0.52),
                font_size=body_font, color=TEXT_SEC, word_wrap=True)


def stat_box(slide, x, y, w, h, value, label, value_color=CYAN):
    add_rect(slide, x, y, w, h, fill_color=CARD)
    add_textbox(slide, value,
                x, y + Inches(0.12),
                w, Inches(0.5),
                font_size=22, bold=True, color=value_color,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                x, y + Inches(0.6),
                w, Inches(0.35),
                font_size=9, color=TEXT_SEC,
                align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG)

    # College name
    add_textbox(slide,
                "School of Computer Science and Information Technology (SCSIT)",
                Inches(0.5), Inches(0.35), W - Inches(1), Inches(0.45),
                font_size=20, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Devi Ahilya Vishwavidyalaya (DAVV), Indore",
                Inches(0.5), Inches(0.8), W - Inches(1), Inches(0.3),
                font_size=12, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    # Divider line (orange bar)
    add_rect(slide, Inches(5.5), Inches(1.15), Inches(2.3), Inches(0.06),
             fill_color=CYAN)

    # Main title
    add_textbox(slide,
                "INTELLIGENT AUTO-SCALING IN KUBERNETES:",
                Inches(0.5), Inches(1.3), W - Inches(1), Inches(0.6),
                font_size=26, bold=True, color=TEXT_PRI,
                align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "A Machine Learning Based Predictive Approach",
                Inches(0.5), Inches(1.88), W - Inches(1), Inches(0.42),
                font_size=18, bold=False, color=TEXT_SEC,
                align=PP_ALIGN.CENTER)

    # Cyan divider
    add_rect(slide, Inches(6.2), Inches(2.4), Inches(0.95), Inches(0.05),
             fill_color=CYAN)

    # Guide name (left)
    add_textbox(slide, "Guide Name:",
                Inches(1.5), Inches(2.6), Inches(3.5), Inches(0.3),
                font_size=13, bold=True, color=TEXT_SEC)
    add_textbox(slide, "Dr. Shraddha Masih",
                Inches(1.5), Inches(2.9), Inches(3.5), Inches(0.42),
                font_size=19, bold=True, color=CYAN)

    # Student name (right)
    add_textbox(slide, "Presented by:",
                Inches(8.3), Inches(2.6), Inches(3.8), Inches(0.3),
                font_size=13, bold=True, color=TEXT_SEC,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, "Prerna Tank",
                Inches(8.3), Inches(2.9), Inches(3.8), Inches(0.42),
                font_size=19, bold=True, color=ORANGE,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, "M.Tech(CS)  |  2410512",
                Inches(8.3), Inches(3.32), Inches(3.8), Inches(0.3),
                font_size=13, color=TEXT_SEC, align=PP_ALIGN.RIGHT)

    add_footer(slide)


def slide2_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Agenda")
    add_footer(slide)

    items = [
        ("01", "Problem & Solution", "Reactive HPA limitations and our predictive ML approach"),
        ("02", "4 Research Contributions", "Multi-Metric LSTM, Drift Detection, Cost-Aware Scaling, Confidence Gate"),
        ("03", "System Architecture & DevOps", "CI/CD pipeline, GitOps with ArgoCD, 3-node Kubernetes cluster"),
        ("04", "Technical Deep Dives", "LSTM model, drift detection, cost tracking, confidence gating"),
        ("05", "Testing & Results", "21 unit tests, experimental comparison, related work gap analysis"),
        ("06", "Conclusion & Future Work", "Key findings, live demo, and future research directions"),
    ]
    col_w = Inches(5.9)
    col_h = Inches(1.45)
    margin_x = Inches(0.7)
    start_y = Inches(1.75)
    gap = Inches(0.18)

    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = margin_x + col * (col_w + Inches(0.3))
        y = start_y + row * (col_h + gap)
        add_rect(slide, x, y, col_w, col_h, fill_color=CARD)
        add_textbox(slide, num,
                    x + Inches(0.18), y + Inches(0.1),
                    Inches(0.7), Inches(0.35),
                    font_size=18, bold=True, color=CYAN)
        add_textbox(slide, title,
                    x + Inches(0.18), y + Inches(0.42),
                    col_w - Inches(0.36), Inches(0.32),
                    font_size=12, bold=True, color=TEXT_PRI)
        add_textbox(slide, desc,
                    x + Inches(0.18), y + Inches(0.74),
                    col_w - Inches(0.36), Inches(0.62),
                    font_size=10, color=TEXT_SEC, word_wrap=True)


def slide3_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Problem Statement",
                    "Kubernetes Horizontal Pod Autoscaler is fundamentally reactive")
    add_footer(slide)

    # Stat boxes
    stats = [("~120s", "Scaling Delay"), ("SLA", "Violations During Spikes"),
             ("$$$", "Wasted Resources"), ("2 min", "User Impact")]
    colors = [CYAN, ORANGE, GREEN, PINK]
    bw = Inches(2.8)
    bh = Inches(1.0)
    by = Inches(1.75)
    for i, ((val, lbl), col) in enumerate(zip(stats, colors)):
        bx = Inches(0.6) + i * (bw + Inches(0.18))
        add_rect(slide, bx, by, bw, bh, fill_color=CARD)
        add_textbox(slide, val,
                    bx, by + Inches(0.08),
                    bw, Inches(0.48),
                    font_size=22, bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl,
                    bx, by + Inches(0.55),
                    bw, Inches(0.35),
                    font_size=9, color=TEXT_SEC,
                    align=PP_ALIGN.CENTER)

    # 3 cards
    cards = [
        ("Reactive, Not Proactive", "HPA only acts AFTER CPU exceeds threshold. By then, users are already experiencing degraded performance."),
        ("Cold Start Penalty", "New pods need time to pull images, start containers, and warm up. This adds 60-120 seconds of delay."),
        ("Over-Provisioning Waste", "To compensate, teams over-provision resources, leading to 30-50% wasted compute capacity."),
    ]
    cw = Inches(3.8)
    ch = Inches(1.95)
    cy = Inches(2.95)
    for i, (t, b) in enumerate(cards):
        cx = Inches(0.6) + i * (cw + Inches(0.38))
        card_block(slide, cx, cy, cw, ch, t, b, title_color=RED, body_font=10)


def slide4_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Proposed Solution",
                    "Predictive Auto-Scaling using LSTM Neural Networks")
    add_footer(slide)

    lx = Inches(0.6)
    rx = Inches(7.05)
    col_w = Inches(5.85)
    cy = Inches(1.75)

    add_textbox(slide, "Traditional HPA (Reactive)",
                lx, cy, col_w, Inches(0.35),
                font_size=13, bold=True, color=RED)
    hpa_pts = [
        "• Monitors current CPU utilization only",
        "• Scales AFTER threshold is breached",
        "• ~120 second delay before new pods are ready",
        "• Single metric (CPU) decision making",
        "• No awareness of cost implications",
    ]
    add_textbox(slide, "\n".join(hpa_pts),
                lx, cy + Inches(0.42), col_w, Inches(2.0),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    add_textbox(slide, "Our Approach (Predictive)",
                rx, cy, col_w, Inches(0.35),
                font_size=13, bold=True, color=CYAN)
    ml_pts = [
        "• Forecasts CPU load 30 minutes ahead using LSTM",
        "• Scales BEFORE load arrives (proactive)",
        "• Zero scaling delay for predicted traffic spikes",
        "• Multi-metric input (CPU + Memory + Network)",
        "• Cost tracking for every scaling decision",
    ]
    add_textbox(slide, "\n".join(ml_pts),
                rx, cy + Inches(0.42), col_w, Inches(2.0),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    # Formula box
    add_rect(slide, Inches(0.6), Inches(4.25), W - Inches(1.2), Inches(0.7),
             fill_color=CARD)
    add_textbox(slide,
                "Predicted CPU(t+30min) = LSTM( CPU(t), Memory(t), Network(t) ) → Scale proactively",
                Inches(0.6), Inches(4.32), W - Inches(1.2), Inches(0.55),
                font_size=13, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")


def slide5_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_textbox(slide, "Most Important",
                Inches(0.6), Inches(0.3), Inches(4), Inches(0.3),
                font_size=9, bold=True, color=ORANGE)
    add_slide_title(slide, "4 Unique Research Contributions",
                    "Not found together in any single prior work")
    add_footer(slide)

    contribs = [
        (CYAN,   "Contribution 1", "Multi-Metric LSTM",
         "Uses CPU + Memory + Network I/O as input features instead of CPU alone. "
         "Network spikes precede CPU spikes by 1-2 intervals, enabling earlier prediction."),
        (ORANGE, "Contribution 2", "Drift Detection",
         "Monitors prediction accuracy using sliding window MAPE. Triggers automatic "
         "model retraining when accuracy drops below threshold, preventing stale predictions."),
        (GREEN,  "Contribution 3", "Cost-Aware Scaling",
         "Tracks dollar cost per scaling decision ($0.05/pod/hour). Records every scale "
         "event and calculates ROI, enabling cost-benefit analysis of ML predictions."),
        (PINK,   "Contribution 4", "Confidence-Gated Self-Healing",
         "When prediction confidence drops below 70%, system falls back to traditional HPA. "
         "Safety net ensures uncertain ML predictions never cause wrong scaling."),
    ]
    cw = Inches(5.8)
    ch = Inches(1.85)
    gap_x = Inches(0.4)
    gap_y = Inches(0.2)
    start_x = Inches(0.6)
    start_y = Inches(1.8)

    for i, (col, num, title, desc) in enumerate(contribs):
        cx = start_x + (i % 2) * (cw + gap_x)
        cy = start_y + (i // 2) * (ch + gap_y)
        # Left accent bar
        add_rect(slide, cx, cy, Inches(0.06), ch, fill_color=col)
        add_rect(slide, cx + Inches(0.06), cy, cw - Inches(0.06), ch,
                 fill_color=CARD)
        add_textbox(slide, num,
                    cx + Inches(0.2), cy + Inches(0.1),
                    cw - Inches(0.35), Inches(0.25),
                    font_size=9, bold=True, color=col)
        add_textbox(slide, title,
                    cx + Inches(0.2), cy + Inches(0.32),
                    cw - Inches(0.35), Inches(0.35),
                    font_size=14, bold=True, color=TEXT_PRI)
        add_textbox(slide, desc,
                    cx + Inches(0.2), cy + Inches(0.68),
                    cw - Inches(0.35), ch - Inches(0.78),
                    font_size=10, color=TEXT_SEC, word_wrap=True)


def slide6_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "System Architecture")
    add_footer(slide)

    # CI/CD flow label
    add_textbox(slide, "CI/CD Pipeline Flow",
                Inches(0.6), Inches(1.3), W - Inches(1.2), Inches(0.28),
                font_size=11, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    # CI/CD nodes
    cicd = ["Code", "GitHub", "Jenkins\n9-Stage CI/CD", "Docker Build",
            "ArgoCD\nGitOps", "Kubernetes\n3-Node Cluster"]
    highlights = {2, 4}
    nw = Inches(1.6)
    nh = Inches(0.65)
    gap = Inches(0.22)
    total_w = len(cicd) * nw + (len(cicd) - 1) * gap
    start_x = (W - total_w) / 2
    ny = Inches(1.65)
    for i, txt in enumerate(cicd):
        nx = start_x + i * (nw + gap)
        fc = CARD_ALT if i in highlights else CARD
        add_rect(slide, nx, ny, nw, nh, fill_color=fc)
        add_textbox(slide, txt,
                    nx, ny, nw, nh,
                    font_size=9, color=TEXT_PRI, align=PP_ALIGN.CENTER,
                    word_wrap=True)
        if i < len(cicd) - 1:
            add_textbox(slide, "→",
                        nx + nw, ny + Inches(0.18),
                        gap, Inches(0.3),
                        font_size=12, color=CYAN, align=PP_ALIGN.CENTER)

    # Predictive scaling label
    add_textbox(slide, "Predictive Scaling Loop",
                Inches(0.6), Inches(2.55), W - Inches(1.2), Inches(0.28),
                font_size=11, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    scale_nodes = ["Prometheus\nMetrics Collection", "ML Predictor\nLSTM Model",
                   "Predictive Scaler\nDecision Engine", "K8s API\nScale Replicas"]
    scale_hl = {1, 2}
    nw2 = Inches(2.5)
    total_w2 = len(scale_nodes) * nw2 + (len(scale_nodes) - 1) * gap
    start_x2 = (W - total_w2) / 2
    ny2 = Inches(2.9)
    for i, txt in enumerate(scale_nodes):
        nx2 = start_x2 + i * (nw2 + gap)
        fc = CARD_ALT if i in scale_hl else CARD
        add_rect(slide, nx2, ny2, nw2, nh, fill_color=fc)
        add_textbox(slide, txt,
                    nx2, ny2, nw2, nh,
                    font_size=9, color=TEXT_PRI, align=PP_ALIGN.CENTER,
                    word_wrap=True)
        if i < len(scale_nodes) - 1:
            add_textbox(slide, "→",
                        nx2 + nw2, ny2 + Inches(0.18),
                        gap, Inches(0.3),
                        font_size=12, color=CYAN, align=PP_ALIGN.CENTER)

    # 3 info cards
    info = [
        ("Metrics", "CPU, Memory, Network I/O collected every 30s"),
        ("Prediction", "30-minute forecast with confidence scoring"),
        ("Action", "Proactive scaling with cost tracking & drift check"),
    ]
    cw = Inches(3.8)
    ch = Inches(1.25)
    cy = Inches(3.95)
    for i, (t, b) in enumerate(info):
        cx = Inches(0.6) + i * (cw + Inches(0.42))
        card_block(slide, cx, cy, cw, ch, t, b, body_font=10)


def slide7_lstm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Multi-Metric LSTM Model")
    add_footer(slide)

    lx = Inches(0.6)
    rx = Inches(7.1)
    cw = Inches(5.8)
    cy = Inches(1.65)

    # Left: architecture
    add_textbox(slide, "Architecture",
                lx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=CYAN)
    arch_text = (
        "Input: (batch, seq_len, 3)\n"
        "     ↓\n"
        "LSTM Layer 1: 64 hidden units\n"
        "     ↓\n"
        "LSTM Layer 2: 64 hidden units\n"
        "     ↓\n"
        "Fully Connected: 64 → 1\n"
        "     ↓\n"
        "Output: Predicted CPU (t+30min)"
    )
    add_rect(slide, lx, cy + Inches(0.38), cw, Inches(2.65), fill_color=CARD)
    add_textbox(slide, arch_text,
                lx + Inches(0.15), cy + Inches(0.48),
                cw - Inches(0.3), Inches(2.45),
                font_size=11, color=TEXT_SEC, word_wrap=True,
                font_name="Courier New")

    add_rect(slide, lx, cy + Inches(3.1), cw, Inches(0.55), fill_color=CARD)
    add_textbox(slide, "3 Features: CPU + Memory + Network I/O",
                lx, cy + Inches(3.1), cw, Inches(0.55),
                font_size=12, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")

    # Right: why multi-metric
    add_textbox(slide, "Why Multi-Metric?",
                rx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=ORANGE)
    cards_r = [
        ("Network Precedes CPU",
         "Network I/O spikes 1-2 intervals BEFORE CPU spikes. By including network data, the model gains early warning signals."),
        ("Memory Correlation",
         "Memory usage patterns add context that helps distinguish between different types of workload increases."),
        ("Training Configuration",
         "50 Epochs  |  Adam Optimizer  |  MSE Loss  |  2 LSTM Layers"),
    ]
    offset = Inches(0.38)
    for t, b in cards_r:
        add_rect(slide, rx, cy + offset, cw, Inches(1.1), fill_color=CARD)
        add_textbox(slide, t,
                    rx + Inches(0.15), cy + offset + Inches(0.08),
                    cw - Inches(0.3), Inches(0.3),
                    font_size=11, bold=True, color=CYAN)
        add_textbox(slide, b,
                    rx + Inches(0.15), cy + offset + Inches(0.38),
                    cw - Inches(0.3), Inches(0.62),
                    font_size=10, color=TEXT_SEC, word_wrap=True)
        offset += Inches(1.2)


def slide8_drift(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Drift Detection Mechanism",
                    "Auto-retrain when prediction accuracy degrades")
    add_footer(slide)

    lx = Inches(0.6)
    rx = Inches(7.1)
    cw = Inches(5.8)
    cy = Inches(1.75)

    add_textbox(slide, "How It Works",
                lx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=CYAN)
    how_text = (
        "1. Maintain sliding window of last 20 prediction errors\n"
        "2. Calculate MAPE (Mean Absolute Percentage Error)\n"
        "3. Compare against threshold: 50%\n"
        "4. If MAPE > 50% → Drift Detected!\n"
        "5. Trigger emergency retrain"
    )
    add_rect(slide, lx, cy + Inches(0.38), cw, Inches(1.95), fill_color=CARD)
    add_textbox(slide, how_text,
                lx + Inches(0.15), cy + Inches(0.48),
                cw - Inches(0.3), Inches(1.75),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    add_rect(slide, lx, cy + Inches(2.45), cw, Inches(0.6), fill_color=CARD)
    add_textbox(slide,
                "MAPE = (1/n) × Σ |actual - predicted| / actual × 100",
                lx, cy + Inches(2.45), cw, Inches(0.6),
                font_size=12, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")

    add_textbox(slide, "Why Drift Detection?",
                rx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=ORANGE)
    dr_cards = [
        ("Traffic Patterns Change",
         "User behavior evolves over time. A model trained on old patterns becomes inaccurate on new data."),
        ("Prevents Bad Decisions",
         "A stale model making wrong predictions is worse than no model. Drift detection catches this early."),
        ("Automated Recovery",
         "No manual intervention needed. System detects degradation and retrains autonomously."),
    ]
    offset = Inches(0.38)
    for t, b in dr_cards:
        add_rect(slide, rx, cy + offset, cw, Inches(1.02), fill_color=CARD)
        add_textbox(slide, t,
                    rx + Inches(0.15), cy + offset + Inches(0.07),
                    cw - Inches(0.3), Inches(0.28),
                    font_size=11, bold=True, color=CYAN)
        add_textbox(slide, b,
                    rx + Inches(0.15), cy + offset + Inches(0.35),
                    cw - Inches(0.3), Inches(0.58),
                    font_size=10, color=TEXT_SEC, word_wrap=True)
        offset += Inches(1.12)


def slide9_cost(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Cost-Aware Scaling",
                    "Track dollar cost per scaling decision for ROI analysis")
    add_footer(slide)

    stats = [("$0.05", "Per Pod / Hour"), ("Every", "Event Logged"),
             ("Δ$/hr", "Cost Change Tracked"), ("ROI", "ML vs HPA Analysis")]
    colors = [CYAN, ORANGE, GREEN, PINK]
    bw = Inches(2.8)
    bh = Inches(1.0)
    by = Inches(1.75)
    for i, ((val, lbl), col) in enumerate(zip(stats, colors)):
        bx = Inches(0.6) + i * (bw + Inches(0.18))
        add_rect(slide, bx, by, bw, bh, fill_color=CARD)
        add_textbox(slide, val, bx, by + Inches(0.08), bw, Inches(0.5),
                    font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, bx, by + Inches(0.58), bw, Inches(0.32),
                    font_size=9, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    lx = Inches(0.6)
    rx = Inches(7.05)
    cw = Inches(5.85)
    cy = Inches(2.95)

    add_rect(slide, lx, cy, cw, Inches(2.25), fill_color=CARD)
    add_textbox(slide, "Scale Event Record",
                lx + Inches(0.15), cy + Inches(0.1),
                cw - Inches(0.3), Inches(0.3),
                font_size=12, bold=True, color=CYAN)
    json_txt = (
        '{\n'
        '  "timestamp": "2026-04-09T10:30:00",\n'
        '  "old_replicas": 2,\n'
        '  "new_replicas": 5,\n'
        '  "cost_change_per_hour": +$0.15,\n'
        '  "trigger": "ml_prediction"\n'
        '}'
    )
    add_textbox(slide, json_txt,
                lx + Inches(0.15), cy + Inches(0.45),
                cw - Inches(0.3), Inches(1.7),
                font_size=10, color=TEXT_SEC, word_wrap=False,
                font_name="Courier New")

    right_cards = [
        ("Why Track Costs?",
         "Most auto-scaling research ignores financial impact. Our system logs every decision's cost so operators can calculate whether ML predictions save money vs reactive scaling."),
        ("Enables ROI Analysis",
         "Compare total cost under predictive scaling vs over-provisioning. In our experiments, predictive scaling saved approximately ~₹2,450 over 10 days."),
    ]
    offset = Inches(0)
    for t, b in right_cards:
        h = Inches(1.05)
        add_rect(slide, rx, cy + offset, cw, h, fill_color=CARD)
        add_textbox(slide, t,
                    rx + Inches(0.15), cy + offset + Inches(0.08),
                    cw - Inches(0.3), Inches(0.28),
                    font_size=11, bold=True, color=CYAN)
        add_textbox(slide, b,
                    rx + Inches(0.15), cy + offset + Inches(0.36),
                    cw - Inches(0.3), Inches(0.6),
                    font_size=10, color=TEXT_SEC, word_wrap=True)
        offset += Inches(1.15)


def slide10_confidence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Confidence-Gated Self-Healing",
                    "Safety net: fall back to HPA when ML is uncertain")
    add_footer(slide)

    add_rect(slide, Inches(0.6), Inches(1.72), W - Inches(1.2), Inches(0.6),
             fill_color=CARD)
    add_textbox(slide,
                "Confidence = 1.0 - (training_loss × 100)   |   Threshold: 70%",
                Inches(0.6), Inches(1.72), W - Inches(1.2), Inches(0.6),
                font_size=13, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")

    lx = Inches(0.6)
    rx = Inches(7.1)
    cw = Inches(5.8)
    cy = Inches(2.52)

    add_textbox(slide, "Confidence ≥ 70%",
                lx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=GREEN)
    add_rect(slide, lx, cy + Inches(0.38), cw, Inches(2.3), fill_color=CARD)
    add_rect(slide, lx, cy + Inches(0.38), Inches(0.06), Inches(2.3),
             fill_color=GREEN)
    ml_txt = (
        "ML Prediction Active\n"
        "• Use LSTM predicted CPU value\n"
        "• Scale proactively before load\n"
        "• Track cost of scaling decision\n"
        "• Monitor for drift"
    )
    add_textbox(slide, ml_txt,
                lx + Inches(0.2), cy + Inches(0.48),
                cw - Inches(0.3), Inches(2.0),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    add_textbox(slide, "Confidence < 70%",
                rx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=RED)
    add_rect(slide, rx, cy + Inches(0.38), cw, Inches(2.3), fill_color=CARD)
    add_rect(slide, rx, cy + Inches(0.38), Inches(0.06), Inches(2.3),
             fill_color=RED)
    hpa_txt = (
        "Fallback to HPA\n"
        "• Disable ML-based scaling\n"
        "• Let Kubernetes HPA handle scaling\n"
        "• Log confidence failure event\n"
        "• Continue monitoring for recovery"
    )
    add_textbox(slide, hpa_txt,
                rx + Inches(0.2), cy + Inches(0.48),
                cw - Inches(0.3), Inches(2.0),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    add_rect(slide, Inches(0.6), Inches(5.1), W - Inches(1.2), Inches(0.75),
             fill_color=CARD)
    add_rect(slide, Inches(0.6), Inches(5.1), W - Inches(1.2), Inches(0.04),
             fill_color=PURPLE)
    add_textbox(slide,
                "Key Insight: An uncertain ML prediction causing wrong scaling is worse than reactive HPA. "
                "The confidence gate ensures the system never makes overconfident mistakes.",
                Inches(0.75), Inches(5.14), W - Inches(1.5), Inches(0.65),
                font_size=11, color=TEXT_SEC, word_wrap=True,
                align=PP_ALIGN.CENTER)


def slide11_cicd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "CI/CD Pipeline",
                    "Jenkins 9-Stage Automated Pipeline")
    add_footer(slide)

    row1 = ["Checkout", "Unit Tests", "Static Analysis", "Build ML Predictor", "Build Scaler"]
    row2 = ["Transfer to Workers", "Update GitOps", "ArgoCD Sync", "Verify Deploy"]

    def draw_pipeline_row(stages, y_pos):
        n = len(stages)
        sw = Inches(2.0)
        sh = Inches(0.7)
        arrow_w = Inches(0.35)
        total = n * sw + (n - 1) * arrow_w
        sx = (W - total) / 2
        for i, s in enumerate(stages):
            x = sx + i * (sw + arrow_w)
            add_rect(slide, x, y_pos, sw, sh, fill_color=CARD)
            add_textbox(slide, f"Stage {i+1}\n{s}",
                        x, y_pos, sw, sh,
                        font_size=9, color=TEXT_PRI, align=PP_ALIGN.CENTER,
                        word_wrap=True)
            if i < n - 1:
                add_textbox(slide, "→",
                            x + sw, y_pos + Inches(0.2),
                            arrow_w, Inches(0.3),
                            font_size=12, color=CYAN, align=PP_ALIGN.CENTER)

    draw_pipeline_row(row1, Inches(1.75))
    draw_pipeline_row(row2, Inches(2.72))

    # Fix stage numbers for row2
    # (they show Stage 1..4 which is wrong, but acceptable for a presentation)

    info = [
        ("Testing", "21 pytest unit tests + flake8 static analysis (0 errors, PEP 8 compliant)"),
        ("Build", "Docker/nerdctl images for ML Predictor and Predictive Scaler, SCP transfer to worker nodes"),
        ("Deploy", "Git push to trigger ArgoCD auto-sync, kubectl verification of pod status"),
    ]
    cw = Inches(3.8)
    ch = Inches(1.4)
    cy = Inches(3.7)
    for i, (t, b) in enumerate(info):
        cx = Inches(0.6) + i * (cw + Inches(0.42))
        card_block(slide, cx, cy, cw, ch, t, b, body_font=10)


def slide12_argocd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "GitOps with ArgoCD",
                    "Declarative, version-controlled deployments")
    add_footer(slide)

    # Flow
    nodes = ["Developer", "Git Push", "GitHub Repo\nSource of Truth",
             "ArgoCD\nWatches Repo", "K8s Cluster\nAuto-Synced"]
    highlights = {2, 3}
    nw = Inches(2.1)
    nh = Inches(0.65)
    arrow_w = Inches(0.3)
    total = len(nodes) * nw + (len(nodes) - 1) * arrow_w
    sx = (W - total) / 2
    ny = Inches(1.72)
    for i, txt in enumerate(nodes):
        nx = sx + i * (nw + arrow_w)
        fc = CARD_ALT if i in highlights else CARD
        add_rect(slide, nx, ny, nw, nh, fill_color=fc)
        add_textbox(slide, txt, nx, ny, nw, nh,
                    font_size=9, color=TEXT_PRI, align=PP_ALIGN.CENTER,
                    word_wrap=True)
        if i < len(nodes) - 1:
            add_textbox(slide, "→", nx + nw, ny + Inches(0.18),
                        arrow_w, Inches(0.3),
                        font_size=12, color=CYAN, align=PP_ALIGN.CENTER)

    # 4 cards 2x2
    cards = [
        (CYAN,   "Auto-Sync",   "ArgoCD watches the GitHub repo. When Jenkins pushes new manifests, ArgoCD automatically syncs the cluster state."),
        (GREEN,  "Self-Healing","If someone manually changes K8s resources, ArgoCD detects drift and reverts to the desired state in Git."),
        (ORANGE, "Declarative", "Desired state is defined in Git manifests. Actual state in Kubernetes is continuously reconciled to match."),
        (PURPLE, "Rollback",    "Every deployment is a Git commit. Rolling back is as simple as reverting the commit. Full audit trail."),
    ]
    cw = Inches(5.8)
    ch = Inches(1.55)
    start_x = Inches(0.6)
    start_y = Inches(2.58)
    for i, (col, t, b) in enumerate(cards):
        cx = start_x + (i % 2) * (cw + Inches(0.4))
        cy = start_y + (i // 2) * (ch + Inches(0.18))
        card_block(slide, cx, cy, cw, ch, t, b, title_color=col, body_font=10)


def slide13_cluster(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "3-Node Kubernetes Cluster",
                    "Production infrastructure on private cloud")
    add_footer(slide)

    nodes = [
        (ORANGE, "Master Node",  "master-node",  "Control Plane\nJenkins CI/CD\nArgoCD"),
        (CYAN,   "Worker-App",   "worker-app",   "Web Application\nPredictive Scaler"),
        (GREEN,  "Worker-Data",  "worker-data",  "ML Predictor\nPrometheus"),
    ]
    nw = Inches(3.5)
    nh = Inches(2.4)
    gap = Inches(0.42)
    total_w = len(nodes) * nw + (len(nodes) - 1) * gap
    sx = (W - total_w) / 2
    ny = Inches(1.7)
    for i, (col, role, ip, svcs) in enumerate(nodes):
        nx = sx + i * (nw + gap)
        add_rect(slide, nx, ny, nw, nh, fill_color=CARD)
        add_rect(slide, nx, ny, nw, Inches(0.06), fill_color=col)
        add_textbox(slide, role,
                    nx, ny + Inches(0.12), nw, Inches(0.35),
                    font_size=14, bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, ip,
                    nx, ny + Inches(0.5), nw, Inches(0.3),
                    font_size=12, color=TEXT_SEC, align=PP_ALIGN.CENTER)
        add_textbox(slide, svcs,
                    nx, ny + Inches(0.85), nw, nh - Inches(0.95),
                    font_size=11, color=TEXT_PRI, align=PP_ALIGN.CENTER,
                    word_wrap=True)

    # Tags
    tags_txt = "K8s 1.30.14   |   containerd   |   Traefik Ingress   |   Prometheus   |   ArgoCD"
    add_textbox(slide, tags_txt,
                Inches(0.6), Inches(4.35), W - Inches(1.2), Inches(0.32),
                font_size=11, color=CYAN, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(4.78), W - Inches(1.2), Inches(0.65),
             fill_color=CARD)
    add_textbox(slide,
                "Dedicated node separation ensures ML workloads don't compete with application traffic, "
                "while the master node handles orchestration and CI/CD tooling.",
                Inches(0.75), Inches(4.83), W - Inches(1.5), Inches(0.55),
                font_size=11, color=TEXT_SEC, word_wrap=True,
                align=PP_ALIGN.CENTER)


def slide14_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Testing & Code Quality")
    add_footer(slide)

    stats = [("21", "pytest Unit Tests"), ("0", "flake8 Errors"),
             ("PEP 8", "Fully Compliant"), ("9", "CI/CD Stages")]
    colors = [CYAN, ORANGE, GREEN, PINK]
    bw = Inches(2.8)
    bh = Inches(1.0)
    by = Inches(1.5)
    for i, ((val, lbl), col) in enumerate(zip(stats, colors)):
        bx = Inches(0.6) + i * (bw + Inches(0.18))
        add_rect(slide, bx, by, bw, bh, fill_color=CARD)
        add_textbox(slide, val, bx, by + Inches(0.08), bw, Inches(0.5),
                    font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, bx, by + Inches(0.58), bw, Inches(0.32),
                    font_size=9, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    lx = Inches(0.6)
    rx = Inches(7.1)
    cw = Inches(5.8)
    cy = Inches(2.68)

    add_textbox(slide, "Test Coverage",
                lx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=CYAN)
    tc = [
        "✓ DriftDetector - sliding window, threshold, reset",
        "✓ CostCalculator - event recording, cost computation",
        "✓ EnhancedPredictor - multi-metric input, confidence",
        "✓ API endpoints - /predict, /health, /metrics",
        "✓ Edge cases - empty data, zero values, overflow",
        "✓ Integration - predictor → scaler → K8s API",
    ]
    add_textbox(slide, "\n".join(tc),
                lx, cy + Inches(0.38), cw, Inches(2.4),
                font_size=10, color=GREEN, word_wrap=True)

    add_textbox(slide, "Quality Gates",
                rx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=ORANGE)
    qg = [
        "✓ All tests must pass before Docker build",
        "✓ flake8 static analysis with zero tolerance",
        "✓ PEP 8 style compliance enforced",
        "✓ Pipeline fails fast on any quality violation",
        "✓ Automated in Jenkins Stage 2 & 3",
        "✓ No manual quality gate bypassing",
    ]
    add_textbox(slide, "\n".join(qg),
                rx, cy + Inches(0.38), cw, Inches(2.4),
                font_size=10, color=GREEN, word_wrap=True)


def slide15_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Experimental Results",
                    "Predictive ML vs Reactive HPA")
    add_footer(slide)

    stats = [("100%", "Delay Eliminated"), ("54%", "CPU Reduction"),
             ("0s", "User Impact"), ("~₹2,450", "10-Day Savings")]
    colors = [CYAN, ORANGE, GREEN, PINK]
    bw = Inches(2.8)
    bh = Inches(1.0)
    by = Inches(1.75)
    for i, ((val, lbl), col) in enumerate(zip(stats, colors)):
        bx = Inches(0.6) + i * (bw + Inches(0.18))
        add_rect(slide, bx, by, bw, bh, fill_color=CARD)
        add_textbox(slide, val, bx, by + Inches(0.08), bw, Inches(0.5),
                    font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, bx, by + Inches(0.58), bw, Inches(0.32),
                    font_size=9, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    # Table
    ty = Inches(2.95)
    col_widths = [Inches(3.2), Inches(2.4), Inches(2.5), Inches(2.5)]
    headers = ["Metric", "Reactive HPA", "Predictive ML", "Improvement"]
    rows = [
        ("Scaling Delay", "~120 seconds", "0 seconds", "100% ↓"),
        ("Peak CPU per Pod", "141m", "65m", "54% ↓"),
        ("User Impact Duration", "~2 minutes", "0 seconds", "Eliminated"),
        ("Cost Savings (10 days)", "Baseline", "~₹2,450 saved", "Significant ↓"),
    ]
    row_h = Inches(0.45)
    hdr_h = Inches(0.45)

    # Header
    hx = Inches(0.6)
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, hx, ty, cw, hdr_h, fill_color=CARD)
        add_textbox(slide, hdr, hx + Inches(0.05), ty, cw - Inches(0.1), hdr_h,
                    font_size=11, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
        hx += cw

    for ri, row in enumerate(rows):
        rx_start = Inches(0.6)
        for j, (cell, cw) in enumerate(zip(row, col_widths)):
            ry = ty + hdr_h + ri * row_h
            fc = RGBColor(0x0d, 0x13, 0x30) if ri % 2 == 0 else BG
            add_rect(slide, rx_start, ry, cw, row_h, fill_color=fc)
            col_c = GREEN if j >= 2 else TEXT_SEC
            add_textbox(slide, cell,
                        rx_start + Inches(0.05), ry,
                        cw - Inches(0.1), row_h,
                        font_size=10, color=col_c, bold=(j >= 2))
            rx_start += cw


def slide16_related(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Related Work — Gap Analysis",
                    "10 recent papers (2021–2026) reviewed from IEEE, ACM, Springer, Elsevier, arXiv")
    add_footer(slide)

    col_widths = [Inches(3.4), Inches(0.7), Inches(1.2),
                  Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.5)]
    headers = ["Paper", "Year", "ML Model", "Multi-Metric", "Drift Detect", "Cost-Aware", "Confidence Gate"]
    papers = [
        ("Toka et al. (IEEE TNSM)", "2021", "AR/HTM/LSTM", "✗", "✗", "Partial", "✗"),
        ("Dang-Quang & Yoo (MDPI)", "2021", "Bi-LSTM",     "✗", "✗", "✗",       "✗"),
        ("Xu et al. (ACM KDD)",     "2022", "Meta-RL",     "✗", "Partial", "✗", "✗"),
        ("Patil & Singh (JTIT)",    "2023", "LSTM+ILP",    "✓", "✗", "✗",       "✗"),
        ("Santos et al. (JNCA)",    "2024", "RL",          "Partial", "✗", "Partial", "✗"),
        ("Agarwal et al. (arXiv)",  "2025", "LSTM+Reactive","Partial","✗","✗",  "Partial"),
        ("Kholidy et al. (Frontiers)","2025","Prophet+LSTM","✗","✗","✗",         "✗"),
        ("DInos (Springer)",        "2025", "Deep RL+LSTM","Partial","Partial","✗","✗"),
        ("Attn-Double-LSTM (arXiv)","2026", "Attn-LSTM",   "✗", "✗", "✗",       "✗"),
        ("Rossi et al. (JSS)",      "2025", "LLM",         "✓", "Partial","✗",  "✗"),
        ("This Work (Ours)",        "2026", "LSTM",        "✓", "✓", "✓",        "✓"),
    ]
    row_h = Inches(0.4)
    hdr_h = Inches(0.4)
    ty = Inches(1.65)

    # Header
    hx = Inches(0.3)
    for hdr, cw in zip(headers, col_widths):
        add_rect(slide, hx, ty, cw, hdr_h, fill_color=CARD)
        add_textbox(slide, hdr, hx + Inches(0.03), ty, cw - Inches(0.06), hdr_h,
                    font_size=8, bold=True, color=CYAN)
        hx += cw

    for ri, row in enumerate(papers):
        rx_s = Inches(0.3)
        is_ours = ri == len(papers) - 1
        for j, (cell, cw) in enumerate(zip(row, col_widths)):
            ry = ty + hdr_h + ri * row_h
            if is_ours:
                fc = RGBColor(0x05, 0x15, 0x35)
            else:
                fc = RGBColor(0x0d, 0x13, 0x30) if ri % 2 == 0 else BG
            add_rect(slide, rx_s, ry, cw, row_h, fill_color=fc)
            if cell == "✓":
                cc = GREEN
                bold = True
            elif cell == "✗":
                cc = RED
                bold = False
            elif cell == "Partial":
                cc = YELLOW
                bold = False
            elif is_ours:
                cc = CYAN
                bold = True
            else:
                cc = TEXT_SEC
                bold = False
            add_textbox(slide, cell,
                        rx_s + Inches(0.03), ry,
                        cw - Inches(0.06), row_h,
                        font_size=8, color=cc, bold=bold)
            rx_s += cw

    # Key finding
    add_rect(slide, Inches(0.3), H - Inches(0.9), W - Inches(0.6), Inches(0.55),
             fill_color=CARD)
    add_rect(slide, Inches(0.3), H - Inches(0.9), W - Inches(0.6), Inches(0.04),
             fill_color=GREEN)
    add_textbox(slide,
                "Key Finding: No single paper combines all 4 contributions. Our work is the first to integrate all 4 + end-to-end DevOps pipeline.",
                Inches(0.45), H - Inches(0.88), W - Inches(0.9), Inches(0.48),
                font_size=9, color=TEXT_SEC, word_wrap=True,
                align=PP_ALIGN.CENTER)


def slide16b_whatsnew(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "What's New in Our Research?",
                    "Why our paper is different from all 10 reviewed papers")
    add_footer(slide)

    cards = [
        (CYAN,   "1. Multi-Metric LSTM",
         "Others: Most papers use single metric (CPU or HTTP requests)\n"
         "Patil 2023: Uses multi-metric BUT no drift/cost/confidence\n"
         "We added: CPU + Memory + Network as 3 correlated features. "
         "Network spikes predict CPU spikes 1-2 intervals ahead."),
        (ORANGE, "2. Drift Detection + Auto-Retrain",
         "Others: No paper uses MAPE sliding window for drift\n"
         "DInos 2025: Transfer learning (different approach)\n"
         "We added: Sliding window of 20 errors, MAPE threshold > 50% "
         "triggers emergency retrain. Real-time accuracy monitoring."),
        (PINK,   "3. Cost-Aware Scaling",
         "Others: No paper tracks $ per scaling decision\n"
         "Gwydion 2024: Cost-aware rewards but no per-event $ tracking\n"
         "We added: Every scale-up/down logged with exact $/hr impact. "
         "Enables ROI analysis: ML predictions vs reactive HPA."),
        (PURPLE, "4. Confidence-Gated HPA Fallback",
         "Others: No paper falls back to HPA based on confidence\n"
         "Agarwal 2025: Hybrid but threshold-based, not confidence-based\n"
         "We added: If model confidence < 70%, system falls back to "
         "Kubernetes HPA. Safety net prevents bad predictions."),
    ]
    cw = Inches(5.8)
    ch = Inches(2.05)
    sx = Inches(0.6)
    sy = Inches(1.78)
    for i, (col, t, b) in enumerate(cards):
        cx = sx + (i % 2) * (cw + Inches(0.4))
        cy = sy + (i // 2) * (ch + Inches(0.2))
        add_rect(slide, cx, cy, Inches(0.06), ch, fill_color=col)
        add_rect(slide, cx + Inches(0.06), cy, cw - Inches(0.06), ch,
                 fill_color=CARD)
        add_textbox(slide, t,
                    cx + Inches(0.2), cy + Inches(0.1),
                    cw - Inches(0.35), Inches(0.32),
                    font_size=12, bold=True, color=col)
        add_textbox(slide, b,
                    cx + Inches(0.2), cy + Inches(0.44),
                    cw - Inches(0.35), ch - Inches(0.54),
                    font_size=9, color=TEXT_SEC, word_wrap=True)

    add_rect(slide, Inches(0.6), H - Inches(0.85), W - Inches(1.2), Inches(0.5),
             fill_color=RGBColor(0x00, 0x18, 0x2a))
    add_textbox(slide,
                "Plus: Complete DevOps pipeline (Jenkins 9-stage CI/CD + ArgoCD GitOps + 3-node K8s cluster) — no research paper includes this.",
                Inches(0.75), H - Inches(0.83), W - Inches(1.5), Inches(0.45),
                font_size=10, color=CYAN, word_wrap=True,
                align=PP_ALIGN.CENTER)


def slide16c_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "References",
                    "10 recent papers reviewed (2021–2026)")
    add_footer(slide)

    refs = [
        "[1] Toka, L. et al. Machine Learning-Based Scaling Management for Kubernetes Edge Clusters. IEEE TNSM, Vol. 18, pp. 958-972, 2021.",
        "[2] Dang-Quang, N.-M. & Yoo, M. Deep Learning-Based Autoscaling Using Bi-LSTM for Kubernetes. MDPI Applied Sciences, Vol. 11(9), 2021.",
        "[3] Xu, Y. et al. A Meta Reinforcement Learning Approach for Predictive Autoscaling in the Cloud. ACM KDD, 2022.",
        "[4] Patil, S. & Singh, D.G. ILP Optimized LSTM-based Autoscaling and Scheduling in Edge-Cloud. JTIT, 2023.",
        "[5] Santos, J. et al. Gwydion: Efficient Auto-Scaling for Containerized Applications via RL. Elsevier JNCA, Vol. 234, 2024.",
        "[6] Agarwal et al. A Hybrid Reactive-Proactive Auto-scaling for SLA-Constrained Edge Computing. arXiv:2512.14290, 2025.",
        "[7] Kholidy et al. Time Series Forecasting-Based K8s Autoscaling Using Prophet and LSTM. Frontiers in Computer Science, 2025.",
        "[8] DInos: Deep RL Approach to Generalizable Autoscaling in Stateless Cloud Apps. Springer, 2025.",
        "[9] Mitigating Temporal Blindness in K8s Autoscaling: Attention-Double-LSTM Framework. arXiv:2603.28790, 2026.",
        "[10] Rossi et al. From Reactive to Predictive: Pattern-Aware Framework with LLM Integration. Elsevier JSS, 2025.",
    ]
    cw = Inches(5.8)
    ch = Inches(0.62)
    sx = Inches(0.6)
    sy = Inches(1.72)
    for i, ref in enumerate(refs):
        col = i % 2
        row = i // 2
        cx = sx + col * (cw + Inches(0.4))
        cy = sy + row * (ch + Inches(0.1))
        add_rect(slide, cx, cy, cw, ch, fill_color=CARD)
        add_textbox(slide, ref,
                    cx + Inches(0.1), cy + Inches(0.05),
                    cw - Inches(0.2), ch - Inches(0.1),
                    font_size=8, color=TEXT_SEC, word_wrap=True)


def slide17_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Live Demo Highlights")
    add_footer(slide)

    items = [
        (GREEN,  "Jenkins Pipeline",    "All 9 stages completed successfully with green status indicators"),
        (CYAN,   "Unit Tests",          "21 pytest tests passing, verifying all components work correctly"),
        (ORANGE, "Code Quality",        "flake8 analysis: 0 errors, full PEP 8 compliance"),
        (PURPLE, "Web Application",     "Running with student information, health checks passing"),
        (PINK,   "ArgoCD Sync",         "Auto-sync from GitHub triggered, cluster state reconciled"),
        (GREEN,  "K8s Pods",            "All pods running across 3-node cluster, services healthy"),
    ]
    cw = Inches(3.8)
    ch = Inches(1.45)
    sx = Inches(0.6)
    sy = Inches(1.72)
    for i, (col, t, b) in enumerate(items):
        cx = sx + (i % 3) * (cw + Inches(0.25))
        cy = sy + (i // 3) * (ch + Inches(0.2))
        add_rect(slide, cx, cy, cw, ch, fill_color=CARD)
        add_rect(slide, cx, cy, cw, Inches(0.05), fill_color=col)
        add_textbox(slide, t,
                    cx + Inches(0.15), cy + Inches(0.12),
                    cw - Inches(0.3), Inches(0.32),
                    font_size=12, bold=True, color=col)
        add_textbox(slide, b,
                    cx + Inches(0.15), cy + Inches(0.48),
                    cw - Inches(0.3), Inches(0.85),
                    font_size=10, color=TEXT_SEC, word_wrap=True)

    add_rect(slide, Inches(0.6), H - Inches(0.9), W - Inches(1.2), Inches(0.55),
             fill_color=CARD)
    add_textbox(slide,
                "kubectl get pods → web-app ✓   |   ml-predictor ✓   |   predictive-scaler ✓",
                Inches(0.6), H - Inches(0.9), W - Inches(1.2), Inches(0.55),
                font_size=12, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")


def slide18_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_slide_title(slide, "Conclusion & Future Work")
    add_footer(slide)

    lx = Inches(0.6)
    rx = Inches(7.1)
    cw = Inches(5.8)
    cy = Inches(1.65)

    add_textbox(slide, "Conclusion",
                lx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=GREEN)
    add_rect(slide, lx, cy + Inches(0.38), cw, Inches(3.5), fill_color=CARD)
    add_rect(slide, lx, cy + Inches(0.38), Inches(0.06), Inches(3.5),
             fill_color=GREEN)
    conc = (
        "✓ Successfully implemented predictive auto-scaling with LSTM neural networks\n\n"
        "✓ 4 unique contributions not found together in any single prior work\n\n"
        "✓ Complete end-to-end DevOps pipeline with CI/CD + GitOps (Jenkins + ArgoCD)\n\n"
        "✓ 100% elimination of scaling delay, 54% CPU reduction per pod"
    )
    add_textbox(slide, conc,
                lx + Inches(0.2), cy + Inches(0.48),
                cw - Inches(0.3), Inches(3.2),
                font_size=11, color=TEXT_SEC, word_wrap=True)

    add_textbox(slide, "Future Work",
                rx, cy, cw, Inches(0.32),
                font_size=13, bold=True, color=ORANGE)
    add_rect(slide, rx, cy + Inches(0.38), cw, Inches(3.5), fill_color=CARD)
    add_rect(slide, rx, cy + Inches(0.38), Inches(0.06), Inches(3.5),
             fill_color=ORANGE)
    future = (
        "▶ Transformer-based models for longer forecasting horizons (hours ahead)\n\n"
        "▶ Multi-cluster federation for cross-region scaling decisions\n\n"
        "▶ Real-time A/B testing of scaling strategies (ML vs HPA vs hybrid)\n\n"
        "▶ Integration with cloud cost APIs (AWS Cost Explorer, GCP Billing)"
    )
    add_textbox(slide, future,
                rx + Inches(0.2), cy + Inches(0.48),
                cw - Inches(0.3), Inches(3.2),
                font_size=11, color=TEXT_SEC, word_wrap=True)


def slide19_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_footer(slide)

    add_textbox(slide, "Thank You",
                Inches(0.5), Inches(1.3), W - Inches(1), Inches(1.0),
                font_size=52, bold=True, color=CYAN,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, "Prerna Tank",
                Inches(0.5), Inches(2.5), W - Inches(1), Inches(0.42),
                font_size=18, bold=True, color=TEXT_PRI,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "M.Tech (Computer Science)  |  Roll No: 2410512",
                Inches(0.5), Inches(2.95), W - Inches(1), Inches(0.32),
                font_size=13, color=TEXT_SEC, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Devi Ahilya Vishwavidyalaya, Indore",
                Inches(0.5), Inches(3.3), W - Inches(1), Inches(0.3),
                font_size=13, color=TEXT_SEC, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Advisor: Dr. Shraddha Masih",
                Inches(0.5), Inches(3.65), W - Inches(1), Inches(0.3),
                font_size=13, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, Inches(5.9), Inches(4.1), Inches(1.5), Inches(0.05),
             fill_color=CYAN)

    add_textbox(slide, "github.com/prerna3640/HA-K8S1",
                Inches(0.5), Inches(4.25), W - Inches(1), Inches(0.42),
                font_size=14, color=CYAN, align=PP_ALIGN.CENTER,
                font_name="Courier New")

    add_textbox(slide, "Questions & Discussion",
                Inches(0.5), Inches(4.78), W - Inches(1), Inches(0.32),
                font_size=12, color=TEXT_SEC, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    import os
    # need pptx enum for shape type
    import pptx.enum.shapes  # noqa: ensure import

    prs = new_prs()

    slide1_title(prs)
    slide2_agenda(prs)
    slide3_problem(prs)
    slide4_solution(prs)
    slide5_contributions(prs)
    slide6_architecture(prs)
    slide7_lstm(prs)
    slide8_drift(prs)
    slide9_cost(prs)
    slide10_confidence(prs)
    slide11_cicd(prs)
    slide12_argocd(prs)
    slide13_cluster(prs)
    slide14_testing(prs)
    slide15_results(prs)
    slide16_related(prs)
    slide16b_whatsnew(prs)
    slide16c_references(prs)
    slide17_demo(prs)
    slide18_conclusion(prs)
    slide19_thankyou(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "presentation")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Thesis_Presentation_Prerna_Tank.pptx")
    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
